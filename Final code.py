from ast import Lambda
from configparser import DEFAULTSECT
from distutils import command
from logging import RootLogger
from statistics import variance
import tkinter as tk
from tkinter import messagebox
from turtle import left
from typing import Literal
from urllib import request
from venv import create
from webbrowser import get
from tkinter import *
from tkinter import filedialog
import moviepy.editor
from moviepy.editor import *
from tkinter.filedialog import *
from PIL import Image
from tkinter.filedialog import askopenfilename,asksaveasfile,askopenfile
from PyPDF2 import PdfFileReader
from numpy import roots
import win32com
import os
import PIL
from PIL import ImageTk,Image
from PIL import Image as PilImage
import pyttsx3
import PyPDF2
from pptx import Presentation
import threading
import sys
from docx2pdf import convert
from tkinter import *
from tkinter import ttk
import os
from random import randint

HEIGHT = 450
WIDTH = 900

''' compress(tar_file,members):
    tar = tarfile.open(tar_file,mode='w:gz')
    progress = tqdm(members)
    for members in progress:
        tar.add(members)
        progress.set_description(f"Compressing{members}")
        tar.close()
'''
'''
#PDF TO WORD
def openFile(): 

    file = askopenfilename(defaultextension=".pdf", 
                                          filetypes=[("Pdf files","*.pdf")])
    if file == "":
        file = None
    else:
        fileEntry.delete(0,END)
        fileEntry.config(fg="blue")
        fileEntry.insert(0,file)

def convert():
    try:
        pdf = fileEntry.get()
        pdfFile = open(pdf, 'rb')
        # creating a pdf reader object
        pdfReader = PdfFileReader(pdfFile) 

        # creating a page object 
        pageObj = pdfReader.getPage(0) 

        # extracting text from page 
        extractedText= pageObj.extractText()
        readPdf.delete(1.0,END)
        readPdf.insert(INSERT,extractedText)

        # closing the pdf file object 
        pdfFile.close()
    except FileNotFoundError:
        fileEntry.delete(0,END)
        fileEntry.config(fg="red")
        fileEntry.insert(0,"Please select a pdf file first")
    except:
        pass




def save2word():
    text = str(readPdf.get(1.0,END))
    wordfile = asksaveasfile(mode='w',defaultextension=".doc", 
                                          filetypes=[("word file",".doc"),
                                                     ("text file",".txt"),
                                                     ("Python file","*.py")])


    if wordfile is None:
        return
    wordfile.write(text)
    wordfile.close()
    print("saved")
    fileEntry.delete(0,END)
    fileEntry.insert(0,"pdf Extracted and Saved...")
#PDF TO WORD
'''
'''
def Word():
    word = win32com.client.Dispatch()
    x = askopenfilename()
    wordObj =  word.Documents.open(x)
    wordObj.SaveAs(word.docx, FileFormat = 16)
'''
'''
def PowerPoint():
    file=askopenfilename(defaultextension=".pptx",
                                filetypes="Ppt files")
    in_file= file
    out_file=os.
'''

class Root(Tk):
    def __init__(self):
        super(Root,self).__init__()
        self.title("Word_to_PDF Converter.")
        self.minsize(300,300)

        #self.wm_iconbitmap('icon.ico')
        #self.config(bg= '#0059b3')
        #self.resizable(width=False, height=False)


        self.lableFrame = ttk.LabelFrame(self, text = "  Open your Word File",relief= "groove")
        self.lableFrame.grid(column = 1, row = 1, padx =20,pady = 20,sticky=N + S + E + W)
        self.button()
        self.make_dir()
        self.button1()

    #define Function for Button1.
    def button(self):
        self.button = ttk.Button(self.lableFrame, text = "Browse a File", command = self.fileDialog)
        #self.img = PhotoImage(file="abc1.png")  # make sure to add "/" not "\"
        #self.button.config(image=self.img)
        self.button.grid(column =1, row = 1)

    # define Function for Button2.
    def button1(self):
        self.button1 = ttk.Button(self.lableFrame, text = "Convert File", command = self.convert)
        #self.img1 = PhotoImage(file="abc1.png")  # make sure to add "/" not "\"
        #self.button1.config(image=self.img1)
        self.button1.grid(column =1, row = 2, padx= 20, pady= 50)

    # define Function for Dialog box for file.
    def fileDialog(self):
        self.filename = filedialog.askopenfilename(initialdir = "/", title = "Select a File", filetype = (("docx", "*.docx"),("All Files", "*.*")))
        self.lable = ttk.Label(self.lableFrame, text = "")
        print(self.filename)
    #define function to create a new folder for output of the file.
    def make_dir(self):
        path = 'D:/Doc_2_PDF (Output)'
        try:
            os.mkdir(path)
        except OSError as error:
            print()
    #define function to convert the docx file into pdf.
    def convert(self):
        i = str(randint(1,1000))
        self.input_file = self.filename
        self.output_file = ('output.pdf')
        convert(self.input_file,self.output_file)


#if __name__ == '__main__':
 #   newwin =Root()
  #  newwin.mainloop()
    
def GUI2():
    window = Root()



def text2speech():
    def playThroughThread():
        global thread
        thread = threading.Thread(target=play, args=())
        thread.daemon = True  # Daemonize thread
        thread.start()  # Start the execution

    def play():
        message='Message: Reading '+file_name_display+' pages from: '+start_pg_entry.get()+' to: '+end_pg_entry.get()
        message_entry.delete(0, END)
        message_entry.insert(0,message)
        document = open(file_name, 'rb')
        speaker = pyttsx3.init()
        text_runs = []
        if var.get() == 1:
            ppt_reader = Presentation(document)
            print(len(ppt_reader.slides))
            for slide in ppt_reader.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        elif var.get()==2:
            pdf_reader = PyPDF2.PdfFileReader(document)
            pages = pdf_reader.numPages
        for page_no in range(pages):
            single_page = pdf_reader.getPage(page_no)
            text_runs.append(single_page.extractText().replace('\n',' '))

        document.close()
        speaker.say(text_runs)
        speaker.runAndWait()



    def stop():
        message='Message: Stopped player'
        message_entry.delete(0, END)
        message_entry.insert(0,message)
        sys.exit()


    def open_file():
        global file_name, file_name_display
        if var.get() == 1:  # PPT
            file = askopenfile(mode='r', filetypes=[('Python Files', '*.pptx')])
            ppt_document = open(file.name, 'rb')
            ppt_reader = Presentation(ppt_document)
            start_pg_entry.config(state='normal')
            start_pg_entry.delete(0,END)
            start_pg_entry.insert(0,1)
            end_pg_entry.config(state='normal')
            end_pg_entry.delete(0,END)
            end_pg_entry.insert(0,len(ppt_reader.slides))
            ppt_document.close()

        elif var.get() == 2:  # PDF
            file = askopenfile(mode='r', filetypes=[('Python Files', '*.pdf')])
            pdf_document = open(file.name, 'rb')
            pdf_reader = PyPDF2.PdfFileReader(pdf_document)
            start_pg_entry.config(state='normal')
            start_pg_entry.delete(0, END)
            start_pg_entry.insert(0, 1)
            end_pg_entry.config(state='normal')
            end_pg_entry.delete(0, END)
            end_pg_entry.insert(0, pdf_reader.numPages)
            pdf_document.close()

        if file is not None:
            file_name = file.name
            print(file.name)
            file_name_display=file_name[file_name.rfind('/')+1: len(file_name)]
            message = 'Message: Choosen file: ' + file_name_display
        else:
            message = 'Message: Error occured while chooing file'
        message_entry.delete(0, END)
        message_entry.insert(0, message)
        


    speech = Toplevel(root)
    speech.title('Read out loud')
    speech.geometry('600x400')

    img = tk.PhotoImage(file='abc1.png')
    img_label= tk.Label(speech,image=img)
    img_label.place(relheight=1,relwidth=1)

    var = IntVar()
#Row1 radio buttons

    #R1= tk.Radiobutton(speech,text='PPT',variable=var,value=1, font=10,bg='black',fg='white')
    #R1.select()
    #R1.place(relx=0.08,rely=0.1,relheight=0.1,relwidth=0.2)
    
    R2= tk.Radiobutton(speech,text='PDF',variable=var,value=2,font=10,bg='black',fg='white')
    R2.place(relx=0.35,rely=0.1,relheight=0.1,relwidth=0.2)
    
    open_file_btn=tk.Button(speech,text='OPEN FILE',font=5,bd=5,bg='black',fg='white',command= lambda: open_file())
    open_file_btn.place(relx=0.6,rely=0.1,relheight=0.1,relwidth=0.35)
    
    start_lbl=tk.Label(speech,text='START PAGE',bg='grey',fg='black',font=1)
    start_lbl.place(relx=0.1,rely=0.28,relheight=0.1,relwidth=0.3)
    
    start_pg_entry=tk.Entry(speech,font=10,bd=5,bg='black',fg='yellow',state='readonly')
    start_pg_entry.place(relx=0.41,rely=0.28,relheight=0.1,relwidth=0.1)
    
    end_lbl=tk.Label(speech,text='END PAGE',font=1,bg='grey',fg='black')
    end_lbl.place(relx=0.52,rely=0.28,relheight=0.1,relwidth=0.3)
    
    end_pg_entry=tk.Entry(speech,bd=5,state='readonly',font=10,bg='black',fg='white')
    end_pg_entry.place(relx=0.83,rely=0.28,relheight=0.1,relwidth=0.1)
    
    message='Message:'
    message_entry=tk.Entry(speech,fg='white',bg='grey',font=10)
    message_entry.place(relx=0.1,rely=0.8,relheight=0.1,relwidth=0.84)
    message_entry.delete(0,END)
    message_entry.insert(0,message)
    
    page_selection_btn=tk.Button(speech,text="PLAY",bd=5,font=10,bg='green',command=playThroughThread)
    page_selection_btn.place(relx=0.15,rely=0.5,relheight=0.2,relwidth=0.35)
    
    stop_btn=tk.Button(speech,text='STOP',bd=5,font=10,bg='red',fg='black',command=stop)
    stop_btn.place(relx=0.52,rely=0.5,relheight=0.2,relwidth=0.35)
    
    
    
    speech.mainloop()

#packing

'''
    img_lbl.grid(row=0, column=0, columnspan=4)
    R1.grid(row=1, column=1 )
    R2.grid(row=1, column=2)
    open_file_btn.grid(row=2, column=0, padx=5, pady=5, columnspan=4)
    start_lbl.grid(row=3,column=0)
    start_pg_entry.grid(row=3,column=1)
    end_lbl.grid(row=3,column=2)
    end_pg_entry.grid(row=3,column=3)
    page_selection_btn.grid(row=4, column=0,  padx=2, columnspan=2,pady=5)
    stop_btn.grid(row=4, column=2, padx=2, columnspan=2,pady=5)
    message_entry.grid(row=10, column=0, columnspan=4, pady=4)
'''
    
    
class ImageCompressor(Tk):
    def __init__(self, winTitle, xSize, ySize, *args):
        super(ImageCompressor, self).__init__()
        if args:
            self.configure(bg=args)
        self.geometry(f'400x250')
        self.title(winTitle)
        self.resizable(False, False)
        self.saveFolder = Button(self,text="Choose which folder to save to", command=self.SavedFolder)
        self.saveFolder.place(x=62.5, y=100)
        self.compressFile = Button(self,text="Choose Image", command=self.GetImageFile)
        self.compressFile.place(x=25, y=15)
        self.chooseQuality = Label(self, text="Choose Image quality", font=("Courier", 10))
        self.chooseQuality.place(x=60, y=70)
        self.scaleValue = Scale(self, from_=100, to=0)
        self.scaleValue.place(x=0, y=70)
        #self.saveFolder = Button(text="Choose which folder to save to", command=self.SavedFolder)
        #self.saveFolder.place(x=62.5, y=100)
        self.imageNameLabel = Label(self,text="Enter new file name")
        self.imageNameLabel.place(x=62.5, y=135)
        self.imageName = Entry(self, bd=3)
        self.imageName.place(x=62.5, y=160)
        self.compressImageBtn = Button(self,text="Compress Image", command=self.CompressImage, bd=5)
        self.compressImageBtn.place(x=270, y=95)
        self.mainloop()

    def GetImageFile(self):
        self.compressLocation = filedialog.askopenfilename()
        if self.compressLocation:
            messagebox.showinfo("File", self.compressLocation)
        else:
            messagebox.showwarning("Error", "No image selected")

    def SavedFolder(self):
        self.saveTo = filedialog.askdirectory()
        if self.saveTo:
            messagebox.showinfo("Save to:", self.saveTo)
        else:
            messagebox.showwarning("Error", "No folder selected")

    def CompressImage(self):
        self.scaleNum = self.scaleValue.get()
        try:
            self.imageToCompress = PilImage.open(self.compressLocation)
            self.getImageExtension = self.compressLocation.rsplit(".", 1)
            self.imageExtension = self.getImageExtension[1]
            self.imageEntryName = self.imageName.get()
            self.imageToCompress.save(f"{self.saveTo}/{self.imageEntryName}.{self.imageExtension}", quality=self.scaleNum)
            messagebox.showinfo("Successful", f"Compressed image saved to {self.saveTo}")
        except:
            messagebox.showwarning("Error", "Something went wrong")

def myGUI():
    MyNewGUI = ImageCompressor("Image Compressor", 450, 225)


    


def openNewWindow():
    
    
    def openFile(): 

        file = askopenfilename(defaultextension=".pdf", 
                                          filetypes=[("Pdf files","*.pdf")])
        if file == "":
            file = None
        else:
            fileEntry.delete(0,END)
            fileEntry.config(fg="blue")
            fileEntry.insert(0,file)

    def convert():
        try: 
            pdf = fileEntry.get()
            pdfFile = open(pdf, 'rb')
        # creating a pdf reader object
            pdfReader = PdfFileReader(pdfFile) 

        # creating a page object 
            pageObj = pdfReader.getPage(0) 

        # extracting text from page 
            extractedText= pageObj.extractText()
            readPdf.delete(1.0,END)
            readPdf.insert(INSERT,extractedText)

        # closing the pdf file object 
            pdfFile.close()

        except FileNotFoundError:
            fileEntry.delete(0,END)
            fileEntry.config(fg="red")
            fileEntry.insert(0,"Please select a pdf file first")
        except:
            pass

    def save2word():
        text = str(readPdf.get(1.0,END))
        wordfile = asksaveasfile(mode='w',defaultextension=".doc", 
                                              filetypes=[("word file",".doc"),
                                                         ("text file",".txt"),
                                                         ("Python file","*.py")])


        if wordfile is None:
            return
        wordfile.write(text)
        wordfile.close()
        print("saved")
        fileEntry.delete(0,END)
        fileEntry.insert(0,"pdf Extracted and Saved...")


    rut = Toplevel(root)
    rut.title('FINAL DOC')
    rut.geometry("600x350")
    rut.config(bg="light blue")


    try:
        rut.wm_iconbitmap("pdf2.ico")
    except:
        print('icon file is not available')
        pass

    
    defaultText = "\n\n\n\n\t\t Your extracted text will apear here.\n \t\t     you can modify that text too."


      


    appName = Label(rut,text="PDF to WORD Converter ",font=('arial',20,'bold'),background='light blue',foreground='maroon')
    appName.place(x=150,y=5)
    #Select pdf file
    labelFile = Label(rut,text="Select Pdf File",font=('arial',12,'bold'))
    labelFile.place(x=30,y=50)
    fileEntry = Entry(rut,font=('calibri',12),fg='white',bg='black',width=40)
    fileEntry.pack(ipadx=200,pady=50,padx=150)

    openFileButton = tk.Button(rut,text=" Open File ",font=('arial',12,'bold'),width=30,
                            background="white",foreground='black',command=openFile)
    openFileButton.place(x=150,y=80)

    convert2Text = tk.Button(rut,text="Read File",font=('arial',12,'bold'),
               background="white",foreground='black',width=15,command=convert)
    convert2Text.place(x=150,y=115)

    readPdf = Text(rut,font=('calibri',12),foreground='light green',background='black',width=60,height=30,bd=10)
    readPdf.pack(padx=20,ipadx=20,pady=20,ipady=20)
    readPdf.insert(INSERT,defaultText)

    save2Word = tk.Button(rut,text="Save to Word File",font=('arial',10,'bold'),width=17,
                    background='white',foreground='black',command=save2word)
    save2Word.place(x=314,y=115)

    rut.mainloop()







def png():
    nav = askopenfilename()
    img=Image.open(nav)
    img.save('xyz'+".jpeg") 

def jpeg():
    monke = askopenfilename()
    img=Image.open(monke)
    img.save('abc'+".png")



'''
VIDEO TO AUDIO
def prntinp():
    inp = entry.get()
    convertedfile(inp);
    '''

def convertedfile():

    video = askopenfilename()
    video = moviepy.editor.VideoFileClip(video)
    audio = video.audio

    audio.write_audiofile("sample.mp3")
    print("Completed!")

#DESKTOP WINDOW 
root=tk.Tk()

root.title("FILE CONVERTER")

canvas= tk.Canvas(root,height=HEIGHT,width=WIDTH)
canvas.pack()

background_image = tk.PhotoImage(file="bg1.png")
background_label = tk.Label(root, image= background_image)
background_label.place(relheight=1,relwidth=1)  

button=tk.Button(root,text='VIDEO to AUDIO',bg='light blue',fg='black',font=20,command=convertedfile)
button.place(relx=0.1,rely=0.1,relheight=0.1,relwidth= 0.41)


button2 = tk.Button(root,text='JPG to PNG',bg='light blue',fg='black',font=20,command = jpeg)
button2.place(relx=0.1,rely=0.22,relheight=0.1,relwidth=0.20)

button4 = tk.Button(root,text='PNG to JPG', bg='light blue',fg='black',font=20,command=png)
button4.place(relx=0.31,rely=0.22,relheight=0.1,relwidth=0.20)


#entry2=tk.Entry(root,bg='black',fg='white',font=20)
#entry2.place(relx=0.5,rely=0.22,relheight=0.1,relwidth= 0.1)


button3= tk.Button(root,text='IMAGE COMPRESSION',bg='light blue',fg='black',font=20,command=myGUI)
button3.place(relx=0.1,rely=0.34,relheight=0.1,relwidth=0.41)


#PDF TO WORD
button5 = tk.Button(root,text='PDF to WORD',bg='light blue',fg='black',font=20,command=openNewWindow)
button5.place(relx=0.1,rely=0.46,relheight=0.1,relwidth=0.41)
'''
fileEntry=tk.Entry(root,bg='black',fg='white',font=20)
fileEntry.place(relx=0.54,rely=0.46,relheight=0.1,relwidth=0.1)

button51 = tk.Button(root,text='convert',bg='white',fg='black',font=20,command=convert)
button51.place(relx=0.65,rely=0.46,relheight=0.1,relwidth=0.14)

button52 = tk.Button(root,text='save file',fg='black',bg='white',font=20,command=save2word)
button52.place(relx=0.8,rely=0.46,relheight=0.1,relwidth=0.14)

readPdf = Text(root,font=('calibri',12),fg='light green',bg='black')
readPdf.place(relx=0.6,rely=0.6,relheight=0.3,relwidth=0.35)
readPdf.insert(INSERT,DEFAULTSECT)
'''
#PDF TO WORD

button6 = tk.Button(root,text='WORD to PDF',bg='light blue',fg='black',font=20,command=GUI2)
button6.place(relx=0.1,rely=0.58,relheight=0.1,relwidth=0.41)

button7 = tk.Button(root,text='TEXT to SPEECH',bg='light blue',fg='black',font=20,command=text2speech)
button7.place(relx=0.1,rely=0.70,relheight=0.1,relwidth=0.41)

root.mainloop() 